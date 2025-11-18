#!/usr/bin/env python3
"""
Commodore C64 vs Amiga - PowerPoint Präsentation Generator v3
Sprites, Betriebssystem, Schnittstellen und Farbtiefe
"""

import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import Rectangle, FancyBboxPatch
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Pfade
MEDIA_DIR = "/home/henry/dock/commodore_Amiga/medien"
OUTPUT_DIR = "/home/henry/dock/commodore_Amiga"

def create_color_palette_comparison(filename):
    """Erstellt Farbpaletten-Vergleich (16 vs 4096 Farben)"""
    fig, axes = plt.subplots(1, 2, figsize=(14, 7))

    # C64 Farbpalette - 16 Farben
    c64_colors = [
        '#000000', '#FFFFFF', '#880000', '#AAFFEE',
        '#CC44CC', '#00CC55', '#0000AA', '#EEEE77',
        '#DD8855', '#664400', '#FF7777', '#333333',
        '#777777', '#AAFF66', '#0088FF', '#BBBBBB'
    ]

    ax1 = axes[0]
    for i, color in enumerate(c64_colors):
        row = i // 4
        col = i % 4
        rect = Rectangle((col * 1.5, (3 - row) * 1.5), 1.3, 1.3,
                         facecolor=color, edgecolor='black', linewidth=2)
        ax1.add_patch(rect)

    ax1.set_xlim(-0.5, 6.5)
    ax1.set_ylim(-0.5, 6.5)
    ax1.set_aspect('equal')
    ax1.axis('off')
    ax1.set_title("C64: 16 Farben\n(VIC-II Chip)", fontsize=16, fontweight='bold', color='#8B4513')

    # Amiga Farbpalette - 4096 Farben (Ausschnitt)
    ax2 = axes[1]

    # Farbverlauf erstellen (4096 = 16x16x16)
    n = 16
    for r in range(n):
        for g in range(n):
            color = f'#{r:01x}{g:01x}8'  # Blau fixiert auf 8
            rect = Rectangle((r * 0.4, g * 0.4), 0.35, 0.35,
                             facecolor=color, edgecolor='none')
            ax2.add_patch(rect)

    ax2.set_xlim(-0.5, 7)
    ax2.set_ylim(-0.5, 7)
    ax2.set_aspect('equal')
    ax2.axis('off')
    ax2.set_title("Amiga: 4096 Farben (HAM)\n(Denise Chip)", fontsize=16, fontweight='bold', color='#4169E1')

    # Verbesserungstext
    fig.text(0.5, 0.02, "+25.500% mehr Farben!", fontsize=16, fontweight='bold',
            color='green', ha='center')

    plt.suptitle("Farbtiefe im Vergleich", fontsize=20, fontweight='bold', y=0.98)
    plt.tight_layout()
    plt.savefig(filename, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Erstellt: {filename}")

def create_sprite_comparison(filename):
    """Erstellt Sprite-Vergleich Visualisierung"""
    fig, axes = plt.subplots(1, 2, figsize=(14, 7))

    # C64 Sprites - 8 Sprites, 24x21 Pixel
    ax1 = axes[0]

    # 8 Sprite-Boxen
    colors = ['#FF0000', '#00FF00', '#0000FF', '#FFFF00',
              '#FF00FF', '#00FFFF', '#FF8800', '#88FF00']

    for i in range(8):
        row = i // 4
        col = i % 4
        # Sprite Box (24x21 proportional)
        rect = Rectangle((col * 2, (1 - row) * 2.5), 1.6, 1.4,
                         facecolor=colors[i], edgecolor='black', linewidth=2, alpha=0.7)
        ax1.add_patch(rect)
        ax1.text(col * 2 + 0.8, (1 - row) * 2.5 + 0.7, f'{i+1}',
                ha='center', va='center', fontsize=12, fontweight='bold')

    ax1.set_xlim(-0.5, 8.5)
    ax1.set_ylim(-0.5, 5)
    ax1.set_aspect('equal')
    ax1.axis('off')
    ax1.set_title("C64: 8 Hardware-Sprites\n24×21 Pixel, 3 Farben", fontsize=14, fontweight='bold', color='#8B4513')

    # Amiga Sprites/BOBs - Unbegrenzt
    ax2 = axes[1]

    # Viele Sprites/BOBs darstellen
    np.random.seed(42)
    for i in range(32):
        x = np.random.uniform(0.5, 7)
        y = np.random.uniform(0.5, 4)
        w = np.random.uniform(0.4, 1.2)
        h = np.random.uniform(0.4, 1.2)
        color = f'#{np.random.randint(0, 256):02x}{np.random.randint(0, 256):02x}{np.random.randint(0, 256):02x}'
        rect = Rectangle((x, y), w, h,
                         facecolor=color, edgecolor='black', linewidth=1, alpha=0.8)
        ax2.add_patch(rect)

    ax2.set_xlim(0, 8.5)
    ax2.set_ylim(0, 5)
    ax2.set_aspect('equal')
    ax2.axis('off')
    ax2.set_title("Amiga: 8 HW-Sprites + BOBs\nUnbegrenzte Größe, 16+ Farben", fontsize=14, fontweight='bold', color='#4169E1')

    plt.suptitle("Sprites & Blitter Objects", fontsize=20, fontweight='bold', y=0.98)
    plt.tight_layout()
    plt.savefig(filename, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Erstellt: {filename}")

def create_interface_comparison(filename):
    """Erstellt Schnittstellen-Vergleich"""
    fig, ax = plt.subplots(figsize=(14, 8))

    # C64 Schnittstellen
    c64_ports = [
        "Expansion Port",
        "RF/TV Out",
        "A/V Port (5-8 Pin)",
        "Serial Port (IEC)",
        "Cassette Port",
        "User Port"
    ]

    # Amiga Schnittstellen
    amiga_ports = [
        "Mouse/Joystick 1",
        "Mouse/Joystick 2",
        "Stereo Audio (RCA)",
        "External Floppy (DB23)",
        "Serial Port (DB25)",
        "Parallel Port (DB25)",
        "RGB Video (DB23)",
        "Composite Video",
        "Expansion Port"
    ]

    # C64 Seite
    for i, port in enumerate(c64_ports):
        y = 7 - i * 1.1
        rect = FancyBboxPatch((0.5, y - 0.4), 4, 0.8, boxstyle="round,pad=0.05",
                              facecolor='#8B4513', edgecolor='black', linewidth=2, alpha=0.7)
        ax.add_patch(rect)
        ax.text(2.5, y, port, ha='center', va='center', fontsize=10, color='white', fontweight='bold')

    ax.text(2.5, 8, "C64\n6 Anschlüsse", ha='center', va='bottom', fontsize=14, fontweight='bold', color='#8B4513')

    # Amiga Seite
    for i, port in enumerate(amiga_ports):
        y = 7.5 - i * 0.85
        rect = FancyBboxPatch((8, y - 0.35), 5, 0.7, boxstyle="round,pad=0.05",
                              facecolor='#4169E1', edgecolor='black', linewidth=2, alpha=0.7)
        ax.add_patch(rect)
        ax.text(10.5, y, port, ha='center', va='center', fontsize=9, color='white', fontweight='bold')

    ax.text(10.5, 8, "Amiga\n9 Anschlüsse", ha='center', va='bottom', fontsize=14, fontweight='bold', color='#4169E1')

    # Verbesserungspfeil
    ax.annotate('', xy=(7.5, 4), xytext=(5.5, 4),
               arrowprops=dict(arrowstyle='->', lw=3, color='green'))
    ax.text(6.5, 4.8, "+50%\nmehr Ports", ha='center', va='bottom',
           fontsize=12, fontweight='bold', color='green')

    ax.set_xlim(0, 14)
    ax.set_ylim(-0.5, 9)
    ax.axis('off')
    ax.set_title("Schnittstellen-Vergleich", fontsize=20, fontweight='bold', y=1.02)

    plt.tight_layout()
    plt.savefig(filename, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Erstellt: {filename}")

def create_os_comparison(filename):
    """Erstellt Betriebssystem-Vergleich"""
    fig, axes = plt.subplots(1, 2, figsize=(14, 7))

    # C64 BASIC V2
    ax1 = axes[0]
    # Blauer Hintergrund
    rect = Rectangle((0, 0), 10, 7, facecolor='#4040E8', edgecolor='black', linewidth=3)
    ax1.add_patch(rect)
    # Rahmen
    rect2 = Rectangle((0.3, 0.3), 9.4, 6.4, facecolor='#4040E8', edgecolor='#7070FF', linewidth=2)
    ax1.add_patch(rect2)

    # Text
    ax1.text(5, 5.5, "**** COMMODORE 64 BASIC V2 ****", ha='center', va='center',
            fontsize=9, color='#7070FF', fontfamily='monospace')
    ax1.text(5, 4.5, "64K RAM SYSTEM  38911 BASIC BYTES FREE", ha='center', va='center',
            fontsize=8, color='#7070FF', fontfamily='monospace')
    ax1.text(5, 3.5, "READY.", ha='center', va='center',
            fontsize=9, color='#7070FF', fontfamily='monospace')
    ax1.text(5, 2.8, "_", ha='center', va='center',
            fontsize=10, color='#7070FF', fontfamily='monospace')

    ax1.set_xlim(0, 10)
    ax1.set_ylim(0, 7)
    ax1.axis('off')
    ax1.set_title("C64: BASIC V2\nKommandozeile", fontsize=14, fontweight='bold', color='#8B4513')

    # Amiga Workbench
    ax2 = axes[1]
    # Grauer Desktop-Hintergrund
    rect = Rectangle((0, 0), 10, 7, facecolor='#AAAAAA', edgecolor='black', linewidth=3)
    ax2.add_patch(rect)

    # Titelleiste
    rect_title = Rectangle((0.2, 5.8), 9.6, 1, facecolor='#0055AA', edgecolor='black', linewidth=1)
    ax2.add_patch(rect_title)
    ax2.text(5, 6.3, "Workbench", ha='center', va='center',
            fontsize=10, color='white', fontweight='bold')

    # Icons
    icons = ["Ram Disk", "Workbench", "Prefs", "Utilities"]
    for i, icon in enumerate(icons):
        x = 1.5 + (i % 2) * 4
        y = 4.5 - (i // 2) * 2
        # Icon Box
        rect_icon = Rectangle((x - 0.6, y - 0.4), 1.2, 0.8, facecolor='#FF8800', edgecolor='black', linewidth=1)
        ax2.add_patch(rect_icon)
        ax2.text(x, y - 0.8, icon, ha='center', va='top', fontsize=7, color='black')

    ax2.set_xlim(0, 10)
    ax2.set_ylim(0, 7)
    ax2.axis('off')
    ax2.set_title("Amiga: Workbench\nGUI mit Multitasking", fontsize=14, fontweight='bold', color='#4169E1')

    plt.suptitle("Betriebssystem-Vergleich", fontsize=20, fontweight='bold', y=0.98)
    plt.tight_layout()
    plt.savefig(filename, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Erstellt: {filename}")

def create_bar_comparison_v2(filename):
    """Erstellt Balkendiagramm für diese Präsentation"""
    categories = ['Sprites\n(Anzahl)', 'Sprite\nFarben', 'System\nFarben', 'Ports']
    c64_values = [8, 3, 16, 6]
    amiga_values = [8, 16, 4096, 9]
    improvements = ['+BOBs', '+433%', '+25500%', '+50%']

    x = np.arange(len(categories))
    width = 0.35

    fig, ax = plt.subplots(figsize=(12, 7))

    bars1 = ax.bar(x - width/2, c64_values, width, label='C64',
                   color='#8B4513', edgecolor='black', linewidth=2)
    bars2 = ax.bar(x + width/2, amiga_values, width, label='Amiga',
                   color='#4169E1', edgecolor='black', linewidth=2)

    ax.set_ylabel('Wert (log. Skala)', fontsize=12, fontweight='bold')
    ax.set_title('Technische Spezifikationen im Vergleich', fontsize=18, fontweight='bold')
    ax.set_xticks(x)
    ax.set_xticklabels(categories, fontsize=11)
    ax.legend(fontsize=12, loc='upper left')
    ax.set_yscale('log')

    # Werte über den Balken
    for bar, val in zip(bars1, c64_values):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() * 1.1,
               str(val), ha='center', va='bottom', fontsize=9, fontweight='bold')

    for bar, val in zip(bars2, amiga_values):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() * 1.1,
               str(val), ha='center', va='bottom', fontsize=9, fontweight='bold')

    # Verbesserungen
    for i, imp in enumerate(improvements):
        ax.text(i, max(amiga_values[i], c64_values[i]) * 3,
               imp, ha='center', va='bottom', fontsize=10,
               fontweight='bold', color='green')

    ax.grid(axis='y', alpha=0.3)
    plt.tight_layout()
    plt.savefig(filename, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Erstellt: {filename}")

def create_presentation():
    """Erstellt die PowerPoint-Präsentation"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[6]  # Blank layout

    # Slide 1: Titelfolie
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Commodore C64 vs. Amiga"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.333), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Sprites, Betriebssystem, Schnittstellen & Farbtiefe"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(102, 102, 102)
    p.alignment = PP_ALIGN.CENTER

    date_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11.333), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Teil 2 - Grafik, System & Konnektivität"
    p.font.size = Pt(20)
    p.font.italic = True
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER

    # Slide 2: Überblick
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Überblick der Verbesserungen"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    img_path = os.path.join(MEDIA_DIR, "bar_comparison_v2.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.3), width=Inches(10.5))

    # Slide 3: Farbtiefe
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Farbtiefe: 16 vs. 4096 Farben"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    img_path = os.path.join(MEDIA_DIR, "color_palette_comparison.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.3), width=Inches(10.5))

    # Slide 4: Echte C64 Farbpalette
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "C64 Farbpalette (VIC-II)"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    c64_palette_path = os.path.join(MEDIA_DIR, "c64_color_palette.png")
    if os.path.exists(c64_palette_path):
        slide.shapes.add_picture(c64_palette_path, Inches(4), Inches(2), width=Inches(5))

    info_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(1.5))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "16 feste Farben - Keine Farbpaletten-Anpassung möglich"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(139, 69, 19)
    p.alignment = PP_ALIGN.CENTER

    # Slide 5: Sprites
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Sprites: Hardware vs. Blitter Objects"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    img_path = os.path.join(MEDIA_DIR, "sprite_comparison.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.3), width=Inches(10.5))

    # Slide 6: Echte C64 Sprites
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "C64 Sprite-Beispiel"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    sprite_path = os.path.join(MEDIA_DIR, "c64_sprites.png")
    if os.path.exists(sprite_path):
        slide.shapes.add_picture(sprite_path, Inches(4), Inches(2), width=Inches(5))

    info_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11), Inches(2))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "8 Hardware-Sprites mit Vergrößerungsmöglichkeit"
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "24×21 Pixel, 3 Farben pro Sprite (+ Hintergrund)"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER

    # Slide 7: Betriebssystem
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Betriebssystem: BASIC V2 vs. AmigaOS"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    img_path = os.path.join(MEDIA_DIR, "os_comparison.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.3), width=Inches(10.5))

    # Slide 8: Echte Workbench
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "AmigaOS Workbench 1.3"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    wb_path = os.path.join(MEDIA_DIR, "amiga_workbench.png")
    if os.path.exists(wb_path):
        slide.shapes.add_picture(wb_path, Inches(2), Inches(1.5), width=Inches(9))

    info_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(1.5))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Preemptives Multitasking - Revolutionär für Heimcomputer!"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.alignment = PP_ALIGN.CENTER

    # Slide 9: Schnittstellen
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Schnittstellen im Vergleich"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    img_path = os.path.join(MEDIA_DIR, "interface_comparison.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.3), width=Inches(10.5))

    # Slide 10: Echte C64 Ports
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "C64 Rückseite - Anschlüsse"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    c64_ports_path = os.path.join(MEDIA_DIR, "c64_rear_ports.gif")
    if os.path.exists(c64_ports_path):
        slide.shapes.add_picture(c64_ports_path, Inches(1.5), Inches(1.8), width=Inches(10))

    info_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(1.5))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Expansion | RF | A/V | Serial (IEC) | Cassette | User Port"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER

    # Slide 11: Echte Amiga Ports
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Amiga 500 Rückseite - Anschlüsse"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    amiga_ports_path = os.path.join(MEDIA_DIR, "amiga_rear_ports.jpg")
    if os.path.exists(amiga_ports_path):
        slide.shapes.add_picture(amiga_ports_path, Inches(1), Inches(2), width=Inches(11))

    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12), Inches(2))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Joystick 1+2 | Stereo Audio | Ext. Floppy | Serial | Parallel | RGB | Composite"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER

    # Slide 12: Zusammenfassung
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Zusammenfassung"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    # Tabelle
    data = [
        ("Eigenschaft", "C64", "Amiga", "Verbesserung"),
        ("Farbtiefe", "16 Farben", "4096 Farben (HAM)", "+25.500%"),
        ("Sprites", "8 HW, 24×21", "8 HW + BOBs", "+Flexibilität"),
        ("Betriebssystem", "BASIC V2", "AmigaOS GUI", "Multitasking"),
        ("Schnittstellen", "6 Ports", "9 Ports", "+50%"),
    ]

    y_pos = 1.5
    for i, row in enumerate(data):
        for j, cell in enumerate(row):
            text_box = slide.shapes.add_textbox(
                Inches(1 + j * 3), Inches(y_pos), Inches(3), Inches(0.6)
            )
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = cell
            p.font.size = Pt(14)
            if i == 0:
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(1 + j * 3), Inches(y_pos - 0.05), Inches(2.9), Inches(0.5)
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(0, 51, 102)
                shape.line.fill.background()
                text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            elif j == 3:
                p.font.color.rgb = RGBColor(0, 128, 0)
                p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
        y_pos += 0.65

    # Fazit
    fazit_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(11), Inches(2))
    tf = fazit_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Die spektakulärste Verbesserung: Farbtiefe mit +25.500%!"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Das preemptive Multitasking von AmigaOS war revolutionär für seine Zeit."
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.CENTER

    # Slide 13: Bildquellen
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Bildquellen"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    sources_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(5))
    tf = sources_box.text_frame
    tf.word_wrap = True

    sources = [
        "Sprite & Farbpaletten-Bilder:",
        "  - c64-wiki.com",
        "",
        "Schnittstellen-Fotos:",
        "  - C64 Ports: c64-wiki.com",
        "  - Amiga Ports: bigbookofamigahardware.com",
        "",
        "Betriebssystem-Screenshots:",
        "  - Amiga Workbench: gregdonner.org",
        "",
        "Visualisierungen: Eigene Erstellung mit Python/Matplotlib"
    ]

    for i, line in enumerate(sources):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        if ":" in line and not line.startswith("  "):
            p.font.bold = True

    # Speichern
    output_path = os.path.join(OUTPUT_DIR, "C64_vs_Amiga_Teil2.pptx")
    prs.save(output_path)
    print(f"\nPräsentation gespeichert: {output_path}")
    return output_path

def main():
    """Hauptfunktion"""
    print("=" * 60)
    print("C64 vs. Amiga - PowerPoint Generator v3")
    print("Sprites, Betriebssystem, Schnittstellen & Farbtiefe")
    print("=" * 60)

    os.makedirs(MEDIA_DIR, exist_ok=True)

    print("\n1. Erstelle Visualisierungen...")

    create_color_palette_comparison(os.path.join(MEDIA_DIR, "color_palette_comparison.png"))
    create_sprite_comparison(os.path.join(MEDIA_DIR, "sprite_comparison.png"))
    create_interface_comparison(os.path.join(MEDIA_DIR, "interface_comparison.png"))
    create_os_comparison(os.path.join(MEDIA_DIR, "os_comparison.png"))
    create_bar_comparison_v2(os.path.join(MEDIA_DIR, "bar_comparison_v2.png"))

    print("\n2. Erstelle PowerPoint-Präsentation...")
    pptx_path = create_presentation()

    print("\n" + "=" * 60)
    print("FERTIG!")
    print(f"Präsentation: {pptx_path}")
    print(f"Medien-Ordner: {MEDIA_DIR}")
    print("=" * 60)

    print("\nDateien im Medien-Ordner:")
    for f in sorted(os.listdir(MEDIA_DIR)):
        size = os.path.getsize(os.path.join(MEDIA_DIR, f))
        print(f"  - {f} ({size/1024:.1f} KB)")

if __name__ == "__main__":
    main()
