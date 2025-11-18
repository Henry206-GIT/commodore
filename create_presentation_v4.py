#!/usr/bin/env python3
"""
Commodore C64 vs Amiga - PowerPoint Präsentation Generator v4
ROM, Grafikauflösung, Sound, Speichermedien, Preis & Verkaufszahlen
"""

import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import Rectangle, FancyBboxPatch
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Pfade
MEDIA_DIR = "/home/henry/dock/commodore_Amiga/medien"
OUTPUT_DIR = "/home/henry/dock/commodore_Amiga"

def create_resolution_comparison(filename):
    """Erstellt Grafikauflösungs-Vergleich"""
    fig, axes = plt.subplots(1, 2, figsize=(14, 7))

    # C64 - 320x200
    ax1 = axes[0]
    rect = Rectangle((0, 0), 320, 200, facecolor='#4040E8', edgecolor='black', linewidth=3)
    ax1.add_patch(rect)

    # Raster-Linien
    for i in range(0, 320, 40):
        ax1.axvline(x=i, color='#6060FF', linewidth=0.5, alpha=0.5)
    for i in range(0, 200, 25):
        ax1.axhline(y=i, color='#6060FF', linewidth=0.5, alpha=0.5)

    ax1.text(160, 100, "320 × 200\n= 64.000 Pixel", ha='center', va='center',
            fontsize=14, color='white', fontweight='bold')

    ax1.set_xlim(-10, 330)
    ax1.set_ylim(-10, 210)
    ax1.set_aspect('equal')
    ax1.axis('off')
    ax1.set_title("C64: 320×200 Pixel\nEinzige Auflösung", fontsize=14, fontweight='bold', color='#8B4513')

    # Amiga - bis zu 1280x800 (skaliert)
    ax2 = axes[1]
    # Hauptauflösung
    rect = Rectangle((0, 0), 320, 256, facecolor='#4169E1', edgecolor='black', linewidth=3, alpha=0.8)
    ax2.add_patch(rect)
    ax2.text(160, 128, "320×256\nStandard", ha='center', va='center',
            fontsize=10, color='white', fontweight='bold')

    # Höhere Auflösungen als Overlay
    rect2 = Rectangle((0, 0), 640, 512, facecolor='none', edgecolor='#00AA00', linewidth=2, linestyle='--')
    ax2.add_patch(rect2)
    ax2.text(320, 480, "640×512 (Hi-Res)", ha='center', va='center',
            fontsize=9, color='#00AA00', fontweight='bold')

    # Interlace
    rect3 = Rectangle((0, 0), 1280, 800, facecolor='none', edgecolor='#FF6600', linewidth=2, linestyle=':')
    ax2.add_patch(rect3)
    ax2.text(640, 750, "1280×800 (Interlace)", ha='center', va='center',
            fontsize=9, color='#FF6600', fontweight='bold')

    ax2.set_xlim(-50, 1350)
    ax2.set_ylim(-50, 850)
    ax2.set_aspect('equal')
    ax2.axis('off')
    ax2.set_title("Amiga: Bis 1280×800 Pixel\nMultiple Auflösungen", fontsize=14, fontweight='bold', color='#4169E1')

    plt.suptitle("Grafikauflösung im Vergleich", fontsize=20, fontweight='bold', y=0.98)
    plt.tight_layout()
    plt.savefig(filename, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Erstellt: {filename}")

def create_sound_comparison(filename):
    """Erstellt Sound-Vergleich (SID vs Paula)"""
    fig, axes = plt.subplots(1, 2, figsize=(14, 7))

    # C64 SID - 3 Voices
    ax1 = axes[0]

    # Drei Synthesizer-Kanäle
    colors = ['#FF6B6B', '#4ECDC4', '#45B7D1']
    labels = ['Voice 1', 'Voice 2', 'Voice 3']

    for i, (color, label) in enumerate(zip(colors, labels)):
        # Wellenform
        x = np.linspace(0, 4*np.pi, 100)
        y = np.sin(x + i*np.pi/3) * (1 - i*0.15) + 2 + i*2.5

        ax1.fill_between(x/(4*np.pi)*8, y - 0.3, y + 0.3, color=color, alpha=0.7)
        ax1.text(9, 2 + i*2.5, label, va='center', fontsize=11, fontweight='bold')

    # Mono-Symbol
    circle = plt.Circle((4, 9), 0.5, color='gray', alpha=0.8)
    ax1.add_patch(circle)
    ax1.text(4, 9, "M", ha='center', va='center', fontsize=10, fontweight='bold', color='white')
    ax1.text(4, 8, "MONO", ha='center', va='top', fontsize=9, color='gray')

    ax1.set_xlim(0, 12)
    ax1.set_ylim(0, 10)
    ax1.axis('off')
    ax1.set_title("C64: SID 6581\n3 Stimmen, Mono, Synthesizer", fontsize=14, fontweight='bold', color='#8B4513')

    # Amiga Paula - 4 Channels Stereo
    ax2 = axes[1]

    # Vier Sample-Kanäle
    colors = ['#FF6B6B', '#4ECDC4', '#45B7D1', '#96CEB4']
    labels = ['Ch 1 (L)', 'Ch 2 (R)', 'Ch 3 (R)', 'Ch 4 (L)']

    for i, (color, label) in enumerate(zip(colors, labels)):
        # Sample-Wellenform (unregelmäßiger)
        x = np.linspace(0, 4*np.pi, 100)
        np.random.seed(i)
        y = np.sin(x + i*np.pi/4) * 0.8 + np.random.randn(100)*0.1 + 1.5 + i*2

        ax2.fill_between(x/(4*np.pi)*8, y - 0.25, y + 0.25, color=color, alpha=0.7)
        ax2.text(9, 1.5 + i*2, label, va='center', fontsize=10, fontweight='bold')

    # Stereo-Symbol
    circle_l = plt.Circle((3, 9.5), 0.4, color='#4169E1', alpha=0.8)
    circle_r = plt.Circle((5, 9.5), 0.4, color='#4169E1', alpha=0.8)
    ax2.add_patch(circle_l)
    ax2.add_patch(circle_r)
    ax2.text(3, 9.5, "L", ha='center', va='center', fontsize=9, fontweight='bold', color='white')
    ax2.text(5, 9.5, "R", ha='center', va='center', fontsize=9, fontweight='bold', color='white')
    ax2.text(4, 8.5, "STEREO", ha='center', va='top', fontsize=9, color='#4169E1')

    ax2.set_xlim(0, 12)
    ax2.set_ylim(0, 10.5)
    ax2.axis('off')
    ax2.set_title("Amiga: Paula 8364\n4 Kanäle, Stereo, 8-Bit Samples", fontsize=14, fontweight='bold', color='#4169E1')

    plt.suptitle("Sound-Hardware im Vergleich", fontsize=20, fontweight='bold', y=0.98)
    plt.tight_layout()
    plt.savefig(filename, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Erstellt: {filename}")

def create_storage_comparison(filename):
    """Erstellt Speichermedien-Vergleich"""
    fig, axes = plt.subplots(1, 2, figsize=(14, 7))

    # C64 - Kassette und 5.25" Diskette
    ax1 = axes[0]

    # Kassette
    rect_cassette = FancyBboxPatch((1, 4), 3, 2, boxstyle="round,pad=0.1",
                                   facecolor='#8B4513', edgecolor='black', linewidth=2)
    ax1.add_patch(rect_cassette)
    # Spulen
    circle1 = plt.Circle((2, 5), 0.4, color='#333', alpha=0.8)
    circle2 = plt.Circle((3, 5), 0.4, color='#333', alpha=0.8)
    ax1.add_patch(circle1)
    ax1.add_patch(circle2)
    ax1.text(2.5, 3.5, "Kassette\n~50 Byte/s", ha='center', va='top', fontsize=10, fontweight='bold')

    # 5.25" Diskette
    rect_floppy = Rectangle((5, 3.5), 3.5, 3.5, facecolor='#333', edgecolor='black', linewidth=2)
    ax1.add_patch(rect_floppy)
    # Loch
    circle3 = plt.Circle((6.75, 5.25), 0.6, color='#666', alpha=0.8)
    ax1.add_patch(circle3)
    # Schlitz
    rect_slot = Rectangle((5.5, 4), 2.5, 0.3, facecolor='#222')
    ax1.add_patch(rect_slot)
    ax1.text(6.75, 3, '5.25" Diskette\n170 KB', ha='center', va='top', fontsize=10, fontweight='bold')

    ax1.set_xlim(0, 10)
    ax1.set_ylim(0, 8)
    ax1.set_aspect('equal')
    ax1.axis('off')
    ax1.set_title("C64: Kassette & 5.25\" Diskette\nLangsam, geringe Kapazität", fontsize=14, fontweight='bold', color='#8B4513')

    # Amiga - 3.5" Diskette
    ax2 = axes[1]

    # 3.5" Diskette (größer dargestellt)
    rect_floppy = Rectangle((2.5, 2), 5, 5, facecolor='#4169E1', edgecolor='black', linewidth=3)
    ax2.add_patch(rect_floppy)

    # Metallschieber
    rect_slider = Rectangle((3.5, 6), 3, 0.8, facecolor='#888', edgecolor='#333', linewidth=1)
    ax2.add_patch(rect_slider)

    # Label-Bereich
    rect_label = Rectangle((3, 2.5), 4, 2, facecolor='white', edgecolor='black', linewidth=1)
    ax2.add_patch(rect_label)
    ax2.text(5, 3.5, "880 KB\nAmiga DOS", ha='center', va='center', fontsize=10, fontweight='bold')

    ax2.text(5, 1.5, '3.5" Diskette\n5x mehr Kapazität', ha='center', va='top', fontsize=11, fontweight='bold')

    ax2.set_xlim(0, 10)
    ax2.set_ylim(0, 8)
    ax2.set_aspect('equal')
    ax2.axis('off')
    ax2.set_title("Amiga: 3.5\" Diskette\nSchneller, zuverlässiger, kompakter", fontsize=14, fontweight='bold', color='#4169E1')

    plt.suptitle("Speichermedien im Vergleich", fontsize=20, fontweight='bold', y=0.98)
    plt.tight_layout()
    plt.savefig(filename, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Erstellt: {filename}")

def create_price_comparison(filename):
    """Erstellt Preis-Vergleich"""
    fig, ax = plt.subplots(figsize=(12, 7))

    # Preise
    products = ['C64', 'Amiga 500']
    prices = [250, 1000]
    colors = ['#8B4513', '#4169E1']

    bars = ax.bar(products, prices, color=colors, edgecolor='black', linewidth=2, width=0.5)

    # Werte über den Balken
    for bar, price in zip(bars, prices):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 30,
               f"${price}", ha='center', va='bottom', fontsize=16, fontweight='bold')

    # Verbesserung
    ax.text(0.5, 600, "+300%\nPreis", ha='center', va='center',
           fontsize=14, fontweight='bold', color='red')

    ax.set_ylabel('Preis (USD)', fontsize=14, fontweight='bold')
    ax.set_title('Preisvergleich bei Markteinführung', fontsize=18, fontweight='bold')
    ax.set_ylim(0, 1200)

    # Zusatzinfo
    ax.text(0, -100, "1982", ha='center', va='top', fontsize=12, color='gray')
    ax.text(1, -100, "1987", ha='center', va='top', fontsize=12, color='gray')

    ax.grid(axis='y', alpha=0.3)
    plt.tight_layout()
    plt.savefig(filename, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Erstellt: {filename}")

def create_sales_comparison(filename):
    """Erstellt Verkaufszahlen-Vergleich"""
    fig, ax = plt.subplots(figsize=(12, 7))

    products = ['C64', 'Amiga\n(alle Modelle)']
    sales = [17, 6]  # In Millionen (C64 hatte eigentlich 12.5-17 Mio)
    colors = ['#8B4513', '#4169E1']

    bars = ax.bar(products, sales, color=colors, edgecolor='black', linewidth=2, width=0.5)

    # Werte über den Balken
    for bar, sale in zip(bars, sales):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.3,
               f"{sale} Mio.", ha='center', va='bottom', fontsize=16, fontweight='bold')

    ax.set_ylabel('Verkaufte Einheiten (Millionen)', fontsize=14, fontweight='bold')
    ax.set_title('Verkaufszahlen im Vergleich', fontsize=18, fontweight='bold')
    ax.set_ylim(0, 20)

    # Info-Text
    ax.text(0.5, 10, "C64:\nMeistverkaufter\nHeimcomputer\naller Zeiten!", ha='center', va='center',
           fontsize=11, fontweight='bold', color='#8B4513',
           bbox=dict(boxstyle='round', facecolor='wheat', alpha=0.8))

    ax.grid(axis='y', alpha=0.3)
    plt.tight_layout()
    plt.savefig(filename, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Erstellt: {filename}")

def create_bar_comparison_v3(filename):
    """Erstellt Balkendiagramm für Teil 3"""
    categories = ['Sound\nKanäle', 'Auflösung\n(max Pixel)', 'Speicher\n(KB)', 'Preis\n(USD)']
    c64_values = [3, 64000, 170, 250]
    amiga_values = [4, 1024000, 880, 1000]
    improvements = ['+33%', '+1500%', '+417%', '+300%']

    x = np.arange(len(categories))
    width = 0.35

    fig, ax = plt.subplots(figsize=(12, 7))

    bars1 = ax.bar(x - width/2, c64_values, width, label='C64',
                   color='#8B4513', edgecolor='black', linewidth=2)
    bars2 = ax.bar(x + width/2, amiga_values, width, label='Amiga',
                   color='#4169E1', edgecolor='black', linewidth=2)

    ax.set_ylabel('Wert (log. Skala)', fontsize=12, fontweight='bold')
    ax.set_title('Technische Spezifikationen im Vergleich', fontsize=18, fontweight='bold')
    ax.set_xticks(x)
    ax.set_xticklabels(categories, fontsize=11)
    ax.legend(fontsize=12, loc='upper left')
    ax.set_yscale('log')

    # Verbesserungen
    for i, imp in enumerate(improvements):
        color = 'green' if i < 3 else 'red'
        ax.text(i, max(amiga_values[i], c64_values[i]) * 2,
               imp, ha='center', va='bottom', fontsize=10,
               fontweight='bold', color=color)

    ax.grid(axis='y', alpha=0.3)
    plt.tight_layout()
    plt.savefig(filename, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Erstellt: {filename}")

def create_presentation():
    """Erstellt die PowerPoint-Präsentation Teil 3"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[6]

    # Slide 1: Titelfolie
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Commodore C64 vs. Amiga"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.333), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Grafikauflösung, Sound, Speichermedien & Preis"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(102, 102, 102)
    p.alignment = PP_ALIGN.CENTER

    date_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11.333), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Teil 3 - Multimedia & Wirtschaftlichkeit"
    p.font.size = Pt(20)
    p.font.italic = True
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER

    # Slide 2: Überblick
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Überblick der Verbesserungen"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    img_path = os.path.join(MEDIA_DIR, "bar_comparison_v3.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.3), width=Inches(10.5))

    # Slide 3: Grafikauflösung
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Grafikauflösung: 320×200 vs. 1280×800"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    img_path = os.path.join(MEDIA_DIR, "resolution_comparison.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.3), width=Inches(10.5))

    # Slide 4: Sound-Vergleich
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Sound: SID vs. Paula"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    img_path = os.path.join(MEDIA_DIR, "sound_comparison.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.3), width=Inches(10.5))

    # Slide 5: Echte SID Chip
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "C64 SID Chip (6581/8580)"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    sid_path = os.path.join(MEDIA_DIR, "sid_chip.jpg")
    if os.path.exists(sid_path):
        slide.shapes.add_picture(sid_path, Inches(3), Inches(1.5), height=Inches(4.5))

    info_box = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11), Inches(1))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "3-stimmiger Synthesizer - Legendär für Chiptunes!"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(139, 69, 19)
    p.alignment = PP_ALIGN.CENTER

    # Slide 6: Echte Paula Chip
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Amiga Paula Chip (8364)"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    paula_path = os.path.join(MEDIA_DIR, "paula_chip.jpg")
    if os.path.exists(paula_path):
        slide.shapes.add_picture(paula_path, Inches(4), Inches(2), width=Inches(5))

    info_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(1.5))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "4 Kanäle, 8-Bit Stereo Samples @ 28 kHz"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(65, 105, 225)
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "Ermöglichte erstmals echte digitale Audiosamples auf Heimcomputern"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER

    # Slide 7: Speichermedien
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Speichermedien im Vergleich"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    img_path = os.path.join(MEDIA_DIR, "storage_comparison.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.3), width=Inches(10.5))

    # Slide 8: C64 Datasette & 1541
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "C64: Datasette & 1541 Floppy"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    # Datasette
    datasette_path = os.path.join(MEDIA_DIR, "datasette.jpg")
    if os.path.exists(datasette_path):
        slide.shapes.add_picture(datasette_path, Inches(1), Inches(1.8), width=Inches(4))

    # 1541
    floppy_path = os.path.join(MEDIA_DIR, "floppy_1541.jpg")
    if os.path.exists(floppy_path):
        slide.shapes.add_picture(floppy_path, Inches(6), Inches(1.8), width=Inches(5.5))

    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(5), Inches(1.5))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Datasette 1530\n~50 Byte/s"
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.CENTER

    info_box2 = slide.shapes.add_textbox(Inches(6), Inches(5.8), Inches(6), Inches(1.5))
    tf = info_box2.text_frame
    p = tf.paragraphs[0]
    p.text = "1541 Floppy Drive\n170 KB, 5.25\""
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.CENTER

    # Slide 9: Amiga Floppy
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Amiga: 3.5\" Floppy Drive"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    amiga_floppy_path = os.path.join(MEDIA_DIR, "amiga_floppy.jpg")
    if os.path.exists(amiga_floppy_path):
        slide.shapes.add_picture(amiga_floppy_path, Inches(3.5), Inches(2), width=Inches(6))

    info_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(1.5))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "880 KB Kapazität - 5x mehr als C64!"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "Kompakter, zuverlässiger, schneller"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER

    # Slide 10: Preis
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Preis: $250 vs. $1000"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    img_path = os.path.join(MEDIA_DIR, "price_comparison.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(2), Inches(1.3), width=Inches(9))

    # Slide 11: Verkaufszahlen
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Verkaufszahlen"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    img_path = os.path.join(MEDIA_DIR, "sales_comparison.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(2), Inches(1.3), width=Inches(9))

    # Slide 12: Zusammenfassung
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Zusammenfassung"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    # Tabelle
    data = [
        ("Eigenschaft", "C64", "Amiga", "Verbesserung"),
        ("Grafikauflösung", "320×200", "bis 1280×800", "+1500%"),
        ("Sound", "SID, 3 Stimmen, Mono", "Paula, 4 Kanäle, Stereo", "+33% Kanäle"),
        ("Speichermedium", "Kassette/5.25\"", "3.5\" Diskette", "+417% Kapazität"),
        ("Preis", "$250", "$1000", "+300%"),
        ("Verkäufe", "~17 Mio.", "~6 Mio.", "C64 führt"),
    ]

    y_pos = 1.4
    for i, row in enumerate(data):
        for j, cell in enumerate(row):
            text_box = slide.shapes.add_textbox(
                Inches(0.8 + j * 3.1), Inches(y_pos), Inches(3), Inches(0.55)
            )
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = cell
            p.font.size = Pt(12)
            if i == 0:
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.8 + j * 3.1), Inches(y_pos - 0.05), Inches(3), Inches(0.45)
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(0, 51, 102)
                shape.line.fill.background()
                text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            elif j == 3:
                p.font.color.rgb = RGBColor(0, 128, 0)
                p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
        y_pos += 0.55

    # Fazit
    fazit_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(11), Inches(2))
    tf = fazit_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Trotz 4x höherem Preis bot der Amiga revolutionäre Multimedia-Fähigkeiten."
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Der C64 bleibt der meistverkaufte Heimcomputer aller Zeiten!"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(139, 69, 19)
    p.alignment = PP_ALIGN.CENTER

    # Slide 13: Bildquellen
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Bildquellen"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    sources_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(5))
    tf = sources_box.text_frame
    tf.word_wrap = True

    sources = [
        "Sound-Chips:",
        "  - SID 6581: c64-wiki.com",
        "  - Paula 8364: bigbookofamigahardware.com",
        "",
        "Speichermedien:",
        "  - Datasette: c64-wiki.com",
        "  - 1541 Floppy: c64-wiki.com",
        "  - Amiga Floppy: bigbookofamigahardware.com",
        "",
        "Visualisierungen: Eigene Erstellung mit Python/Matplotlib"
    ]

    for i, line in enumerate(sources):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        if ":" in line and not line.startswith("  "):
            p.font.bold = True

    # Speichern
    output_path = os.path.join(OUTPUT_DIR, "C64_vs_Amiga_Teil3.pptx")
    prs.save(output_path)
    print(f"\nPräsentation gespeichert: {output_path}")
    return output_path

def main():
    """Hauptfunktion"""
    print("=" * 60)
    print("C64 vs. Amiga - PowerPoint Generator v4")
    print("Grafikauflösung, Sound, Speichermedien & Preis")
    print("=" * 60)

    os.makedirs(MEDIA_DIR, exist_ok=True)

    print("\n1. Erstelle Visualisierungen...")

    create_resolution_comparison(os.path.join(MEDIA_DIR, "resolution_comparison.png"))
    create_sound_comparison(os.path.join(MEDIA_DIR, "sound_comparison.png"))
    create_storage_comparison(os.path.join(MEDIA_DIR, "storage_comparison.png"))
    create_price_comparison(os.path.join(MEDIA_DIR, "price_comparison.png"))
    create_sales_comparison(os.path.join(MEDIA_DIR, "sales_comparison.png"))
    create_bar_comparison_v3(os.path.join(MEDIA_DIR, "bar_comparison_v3.png"))

    print("\n2. Erstelle PowerPoint-Präsentation...")
    pptx_path = create_presentation()

    print("\n" + "=" * 60)
    print("FERTIG!")
    print(f"Präsentation: {pptx_path}")
    print(f"Medien-Ordner: {MEDIA_DIR}")
    print("=" * 60)

    print("\nDateien im Medien-Ordner:")
    for f in sorted(os.listdir(MEDIA_DIR)):
        size = os.path.getsize(os.path.join(MEDIA_DIR, f))
        print(f"  - {f} ({size/1024:.1f} KB)")

if __name__ == "__main__":
    main()
