#!/usr/bin/env python3
"""
Commodore C64 vs Amiga - PowerPoint Präsentation Generator
Erstellt Visualisierungen und eine PowerPoint-Präsentation
"""

import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch, Rectangle, Wedge
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Pfade
MEDIA_DIR = "/home/henry/dock/commodore_Amiga/medien"
OUTPUT_DIR = "/home/henry/dock/commodore_Amiga"

def create_speedometer(filename):
    """Erstellt Tachometer-Vergleich für Taktfrequenz (1 MHz vs 7 MHz)"""
    fig, axes = plt.subplots(1, 2, figsize=(14, 6))

    for ax, (name, value, max_val, color) in zip(axes, [
        ("C64", 1, 10, "#8B4513"),
        ("Amiga", 7, 10, "#4169E1")
    ]):
        # Tachometer-Hintergrund
        theta = np.linspace(np.pi, 0, 100)
        ax.fill_between(theta, 0.6, 1.0, alpha=0.3, color='lightgray', transform=ax.transData)

        # Skala zeichnen
        for i in range(11):
            angle = np.pi - (i * np.pi / 10)
            x_outer = np.cos(angle)
            y_outer = np.sin(angle)
            x_inner = 0.85 * np.cos(angle)
            y_inner = 0.85 * np.sin(angle)
            ax.plot([x_inner, x_outer], [y_inner, y_outer], 'k-', linewidth=2)

            # Beschriftung
            x_text = 0.7 * np.cos(angle)
            y_text = 0.7 * np.sin(angle)
            ax.text(x_text, y_text, str(i), ha='center', va='center', fontsize=10, fontweight='bold')

        # Wertanzeige
        value_angle = np.pi - (value * np.pi / max_val)
        ax.arrow(0, 0, 0.6 * np.cos(value_angle), 0.6 * np.sin(value_angle),
                head_width=0.08, head_length=0.05, fc=color, ec=color, linewidth=3)

        # Kreismitte
        circle = plt.Circle((0, 0), 0.1, color=color, zorder=5)
        ax.add_patch(circle)

        # Halbbogen
        arc = Wedge((0, 0), 1.0, 0, 180, width=0.15, facecolor=color, alpha=0.6)
        ax.add_patch(arc)

        ax.set_xlim(-1.3, 1.3)
        ax.set_ylim(-0.3, 1.3)
        ax.set_aspect('equal')
        ax.axis('off')
        ax.set_title(f"{name}\n{value} MHz", fontsize=18, fontweight='bold', color=color)

    plt.suptitle("Taktfrequenz-Vergleich", fontsize=22, fontweight='bold', y=1.02)
    plt.tight_layout()
    plt.savefig(filename, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Erstellt: {filename}")

def create_water_glasses(filename):
    """Erstellt Wasserglas-Visualisierung für RAM (64 KB vs 512 KB)"""
    fig, ax = plt.subplots(figsize=(12, 8))

    # Glas 1 - C64 (64 KB)
    glass1_x, glass1_y = 2, 0
    glass1_width, glass1_height = 2, 6

    # Glasform C64
    glass1 = FancyBboxPatch((glass1_x, glass1_y), glass1_width, glass1_height,
                            boxstyle="round,pad=0.05", facecolor='lightcyan',
                            edgecolor='black', linewidth=3, alpha=0.5)
    ax.add_patch(glass1)

    # Wasserfüllstand C64 (64/512 = 12.5%)
    fill_height1 = glass1_height * (64 / 512)
    water1 = Rectangle((glass1_x + 0.1, glass1_y + 0.1),
                       glass1_width - 0.2, fill_height1,
                       facecolor='#8B4513', alpha=0.8)
    ax.add_patch(water1)

    # Glas 2 - Amiga (512 KB)
    glass2_x, glass2_y = 7, 0
    glass2_width, glass2_height = 2, 6

    # Glasform Amiga
    glass2 = FancyBboxPatch((glass2_x, glass2_y), glass2_width, glass2_height,
                            boxstyle="round,pad=0.05", facecolor='lightcyan',
                            edgecolor='black', linewidth=3, alpha=0.5)
    ax.add_patch(glass2)

    # Wasserfüllstand Amiga (512/512 = 100%)
    fill_height2 = glass2_height * (512 / 512)
    water2 = Rectangle((glass2_x + 0.1, glass2_y + 0.1),
                       glass2_width - 0.2, fill_height2 - 0.2,
                       facecolor='#4169E1', alpha=0.8)
    ax.add_patch(water2)

    # Beschriftungen
    ax.text(glass1_x + glass1_width/2, glass1_height + 0.8,
           "C64\n64 KB", ha='center', va='bottom',
           fontsize=16, fontweight='bold', color='#8B4513')

    ax.text(glass2_x + glass2_width/2, glass2_height + 0.8,
           "Amiga\n512 KB", ha='center', va='bottom',
           fontsize=16, fontweight='bold', color='#4169E1')

    # Vergleichspfeil
    ax.annotate('', xy=(6.5, 3), xytext=(4.5, 3),
               arrowprops=dict(arrowstyle='->', lw=3, color='green'))
    ax.text(5.5, 3.5, "+700%", ha='center', va='bottom',
           fontsize=14, fontweight='bold', color='green')

    ax.set_xlim(0, 11)
    ax.set_ylim(-1, 9)
    ax.set_aspect('equal')
    ax.axis('off')
    ax.set_title("Arbeitsspeicher-Vergleich (RAM)", fontsize=20, fontweight='bold', y=1.05)

    plt.tight_layout()
    plt.savefig(filename, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Erstellt: {filename}")

def create_databus_visualization(filename):
    """Erstellt 8-Bit vs 16-Bit Datenbus-Visualisierung"""
    fig, axes = plt.subplots(1, 2, figsize=(14, 6))

    # C64 - 8-Bit Bus
    ax1 = axes[0]
    for i in range(8):
        rect = Rectangle((0.5, i * 0.8), 4, 0.6,
                         facecolor='#8B4513', edgecolor='black', linewidth=2)
        ax1.add_patch(rect)
        ax1.text(2.5, i * 0.8 + 0.3, f"Bit {i}", ha='center', va='center',
                fontsize=10, color='white', fontweight='bold')

    # Datenfluss-Pfeile
    for i in range(4):
        ax1.annotate('', xy=(5.5, i * 2 + 0.5), xytext=(4.7, i * 2 + 0.5),
                    arrowprops=dict(arrowstyle='->', lw=2, color='#8B4513'))

    ax1.set_xlim(0, 7)
    ax1.set_ylim(-1, 7)
    ax1.axis('off')
    ax1.set_title("C64: 8-Bit Datenbus\nMOS 6510", fontsize=16, fontweight='bold', color='#8B4513')

    # Amiga - 16-Bit Bus
    ax2 = axes[1]
    for i in range(16):
        rect = Rectangle((0.5, i * 0.4), 4, 0.3,
                         facecolor='#4169E1', edgecolor='black', linewidth=1.5)
        ax2.add_patch(rect)
        if i % 2 == 0:
            ax2.text(2.5, i * 0.4 + 0.15, f"{i}", ha='center', va='center',
                    fontsize=7, color='white', fontweight='bold')

    # Datenfluss-Pfeile (doppelt so viele)
    for i in range(8):
        ax2.annotate('', xy=(5.5, i * 0.8 + 0.2), xytext=(4.7, i * 0.8 + 0.2),
                    arrowprops=dict(arrowstyle='->', lw=2, color='#4169E1'))

    ax2.set_xlim(0, 7)
    ax2.set_ylim(-1, 7)
    ax2.axis('off')
    ax2.set_title("Amiga: 16-Bit Datenbus\nMotorola 68000", fontsize=16, fontweight='bold', color='#4169E1')

    plt.suptitle("Datenbus-Breite Vergleich", fontsize=20, fontweight='bold', y=1.02)
    plt.tight_layout()
    plt.savefig(filename, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Erstellt: {filename}")

def create_memory_blocks(filename):
    """Erstellt Stapel von 1 KB-Blöcken (64 vs 512)"""
    fig, axes = plt.subplots(1, 2, figsize=(14, 8))

    # C64 - 64 KB (8x8 Grid)
    ax1 = axes[0]
    blocks_c64 = 64
    cols_c64 = 8
    rows_c64 = 8

    for i in range(blocks_c64):
        row = i // cols_c64
        col = i % cols_c64
        rect = Rectangle((col * 1.1, row * 1.1), 1, 1,
                         facecolor='#8B4513', edgecolor='black', linewidth=1, alpha=0.8)
        ax1.add_patch(rect)

    ax1.set_xlim(-0.5, cols_c64 * 1.1 + 0.5)
    ax1.set_ylim(-0.5, rows_c64 * 1.1 + 0.5)
    ax1.set_aspect('equal')
    ax1.axis('off')
    ax1.set_title(f"C64: {blocks_c64} KB\n({rows_c64}x{cols_c64} Blöcke à 1 KB)",
                 fontsize=16, fontweight='bold', color='#8B4513')

    # Amiga - 512 KB (16x32 Grid, skaliert dargestellt)
    ax2 = axes[1]
    blocks_amiga = 512
    cols_amiga = 32
    rows_amiga = 16

    for i in range(blocks_amiga):
        row = i // cols_amiga
        col = i % cols_amiga
        rect = Rectangle((col * 0.28, row * 0.55), 0.25, 0.5,
                         facecolor='#4169E1', edgecolor='black', linewidth=0.5, alpha=0.8)
        ax2.add_patch(rect)

    ax2.set_xlim(-0.5, cols_amiga * 0.28 + 0.5)
    ax2.set_ylim(-0.5, rows_amiga * 0.55 + 0.5)
    ax2.set_aspect('equal')
    ax2.axis('off')
    ax2.set_title(f"Amiga: {blocks_amiga} KB\n({rows_amiga}x{cols_amiga} Blöcke à 1 KB)",
                 fontsize=16, fontweight='bold', color='#4169E1')

    plt.suptitle("Speicherblock-Vergleich", fontsize=20, fontweight='bold', y=1.02)
    plt.tight_layout()
    plt.savefig(filename, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Erstellt: {filename}")

def create_processor_comparison(filename):
    """Erstellt Prozessor-Vergleichsgrafik"""
    fig, ax = plt.subplots(figsize=(12, 8))

    # C64 Prozessor
    c64_chip = FancyBboxPatch((1, 3), 4, 3, boxstyle="round,pad=0.1",
                              facecolor='#8B4513', edgecolor='black', linewidth=3)
    ax.add_patch(c64_chip)

    # Pins für C64
    for i in range(8):
        ax.plot([0.5, 1], [3.3 + i * 0.3, 3.3 + i * 0.3], 'k-', linewidth=2)
        ax.plot([5, 5.5], [3.3 + i * 0.3, 3.3 + i * 0.3], 'k-', linewidth=2)

    ax.text(3, 4.5, "MOS 6510", ha='center', va='center',
           fontsize=14, fontweight='bold', color='white')
    ax.text(3, 4, "8-Bit", ha='center', va='center',
           fontsize=12, color='white')
    ax.text(3, 3.5, "1 MHz", ha='center', va='center',
           fontsize=11, color='lightgray')

    # Amiga Prozessor (größer)
    amiga_chip = FancyBboxPatch((7, 2), 5, 5, boxstyle="round,pad=0.1",
                                facecolor='#4169E1', edgecolor='black', linewidth=3)
    ax.add_patch(amiga_chip)

    # Pins für Amiga (mehr Pins)
    for i in range(16):
        ax.plot([6.5, 7], [2.3 + i * 0.28, 2.3 + i * 0.28], 'k-', linewidth=1.5)
        ax.plot([12, 12.5], [2.3 + i * 0.28, 2.3 + i * 0.28], 'k-', linewidth=1.5)

    ax.text(9.5, 4.8, "Motorola", ha='center', va='center',
           fontsize=14, fontweight='bold', color='white')
    ax.text(9.5, 4.2, "68000", ha='center', va='center',
           fontsize=14, fontweight='bold', color='white')
    ax.text(9.5, 3.5, "16-Bit", ha='center', va='center',
           fontsize=12, color='white')
    ax.text(9.5, 2.9, "7 MHz", ha='center', va='center',
           fontsize=11, color='lightgray')

    # Verbesserungspfeile
    ax.annotate('', xy=(6.8, 4.5), xytext=(5.7, 4.5),
               arrowprops=dict(arrowstyle='->', lw=3, color='green'))
    ax.text(6.25, 5.2, "+700%\nArchitektur", ha='center', va='bottom',
           fontsize=10, fontweight='bold', color='green')
    ax.text(6.25, 3.8, "+600%\nTakt", ha='center', va='top',
           fontsize=10, fontweight='bold', color='green')

    ax.set_xlim(0, 13)
    ax.set_ylim(1, 8)
    ax.axis('off')
    ax.set_title("Prozessor-Vergleich: MOS 6510 vs Motorola 68000",
                fontsize=18, fontweight='bold', y=1.02)

    plt.tight_layout()
    plt.savefig(filename, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Erstellt: {filename}")

def create_bar_comparison(filename):
    """Erstellt Balkendiagramm für alle Verbesserungen"""
    categories = ['Prozessor\n(Bit-Breite)', 'Taktfrequenz', 'RAM', 'Farbtiefe']
    c64_values = [8, 1, 64, 16]
    amiga_values = [16, 7, 512, 4096]
    improvements = ['+100%', '+600%', '+700%', '+25500%']

    x = np.arange(len(categories))
    width = 0.35

    fig, ax = plt.subplots(figsize=(12, 7))

    # Logarithmische Skala für bessere Darstellung
    bars1 = ax.bar(x - width/2, c64_values, width, label='C64',
                   color='#8B4513', edgecolor='black', linewidth=2)
    bars2 = ax.bar(x + width/2, amiga_values, width, label='Amiga',
                   color='#4169E1', edgecolor='black', linewidth=2)

    ax.set_ylabel('Wert (log. Skala)', fontsize=12, fontweight='bold')
    ax.set_title('Technische Spezifikationen im Vergleich', fontsize=18, fontweight='bold')
    ax.set_xticks(x)
    ax.set_xticklabels(categories, fontsize=11)
    ax.legend(fontsize=12, loc='upper left')
    ax.set_yscale('log')

    # Werte über den Balken
    for bar, val in zip(bars1, c64_values):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() * 1.1,
               str(val), ha='center', va='bottom', fontsize=9, fontweight='bold')

    for bar, val in zip(bars2, amiga_values):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() * 1.1,
               str(val), ha='center', va='bottom', fontsize=9, fontweight='bold')

    # Verbesserungen anzeigen
    for i, imp in enumerate(improvements):
        ax.text(i, max(amiga_values[i], c64_values[i]) * 3,
               imp, ha='center', va='bottom', fontsize=10,
               fontweight='bold', color='green')

    ax.grid(axis='y', alpha=0.3)
    plt.tight_layout()
    plt.savefig(filename, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Erstellt: {filename}")

def create_presentation():
    """Erstellt die PowerPoint-Präsentation"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Titelfolie
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Titel
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Commodore C64 vs. Amiga"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # Untertitel
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.333), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Technischer Vergleich: Prozessor, Taktfrequenz & Arbeitsspeicher"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(102, 102, 102)
    p.alignment = PP_ALIGN.CENTER

    # Zeitraum
    date_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11.333), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "1982 vs. 1985 - Der Sprung in eine neue Ära"
    p.font.size = Pt(20)
    p.font.italic = True
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER

    # Slide 2: Überblick
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Überblick der Verbesserungen"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    # Bild einfügen
    img_path = os.path.join(MEDIA_DIR, "bar_comparison.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.3), width=Inches(10.5))

    # Slide 3: Prozessor-Vergleich
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Prozessor: MOS 6510 vs. Motorola 68000"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    img_path = os.path.join(MEDIA_DIR, "processor_comparison.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.5), width=Inches(10))

    # Slide 4: Datenbus-Visualisierung
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Datenbus: 8-Bit vs. 16-Bit Architektur"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    img_path = os.path.join(MEDIA_DIR, "databus_visualization.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.3), width=Inches(10.5))

    # Bullet Points
    bullet_box = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11), Inches(1))
    tf = bullet_box.text_frame
    p = tf.paragraphs[0]
    p.text = "16-Bit = Doppelte Datenbreite pro Taktzyklus = Mehr Leistung"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER

    # Slide 5: Taktfrequenz
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Taktfrequenz: 1 MHz vs. 7 MHz"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    img_path = os.path.join(MEDIA_DIR, "speedometer.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.3), width=Inches(10.5))

    # Info-Box
    info_box = slide.shapes.add_textbox(Inches(2), Inches(6.2), Inches(9), Inches(0.8))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "+600% Geschwindigkeitssteigerung - 7x mehr Operationen pro Sekunde"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.alignment = PP_ALIGN.CENTER

    # Slide 6: RAM als Wassergläser
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Arbeitsspeicher: 64 KB vs. 512 KB"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    img_path = os.path.join(MEDIA_DIR, "water_glasses.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(2), Inches(1.3), width=Inches(9))

    # Slide 7: RAM als Blöcke
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "RAM-Blöcke: 64 vs. 512 Kilobyte"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    img_path = os.path.join(MEDIA_DIR, "memory_blocks.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.3), width=Inches(10.5))

    info_box = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11), Inches(0.8))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Jeder Block = 1 KB | 8x mehr Speicher für Programme und Daten"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER

    # Slide 8: Zusammenfassung
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Zusammenfassung"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    # Zusammenfassungs-Tabelle
    data = [
        ("Eigenschaft", "C64", "Amiga", "Verbesserung"),
        ("Prozessor", "MOS 6510 (8-Bit)", "Motorola 68000 (16-Bit)", "+700%"),
        ("Taktfrequenz", "1 MHz", "7 MHz", "+600%"),
        ("RAM", "64 KB", "512 KB", "+700%"),
    ]

    # Einfache Textdarstellung der Tabelle
    y_pos = 1.5
    for i, row in enumerate(data):
        for j, cell in enumerate(row):
            text_box = slide.shapes.add_textbox(
                Inches(1 + j * 3), Inches(y_pos), Inches(3), Inches(0.6)
            )
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = cell
            p.font.size = Pt(16)
            if i == 0:
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                # Header-Hintergrund
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(1 + j * 3), Inches(y_pos - 0.05), Inches(2.9), Inches(0.5)
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(0, 51, 102)
                shape.line.fill.background()
                # Text nach vorne bringen
                text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            elif j == 3:  # Verbesserungsspalte
                p.font.color.rgb = RGBColor(0, 128, 0)
                p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
        y_pos += 0.7

    # Fazit
    fazit_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11), Inches(2))
    tf = fazit_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Der Amiga war ein technisches Wunder seiner Zeit und übertraf seinen Vorgänger"
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "sowie die Konkurrenz massiv mit seinen Verbesserungen."
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER

    # Speichern
    output_path = os.path.join(OUTPUT_DIR, "C64_vs_Amiga_Vergleich.pptx")
    prs.save(output_path)
    print(f"\nPräsentation gespeichert: {output_path}")
    return output_path

def main():
    """Hauptfunktion"""
    print("=" * 60)
    print("C64 vs. Amiga - PowerPoint Generator")
    print("=" * 60)

    # Sicherstellen, dass der Medien-Ordner existiert
    os.makedirs(MEDIA_DIR, exist_ok=True)

    print("\n1. Erstelle Visualisierungen...")

    # Visualisierungen erstellen
    create_speedometer(os.path.join(MEDIA_DIR, "speedometer.png"))
    create_water_glasses(os.path.join(MEDIA_DIR, "water_glasses.png"))
    create_databus_visualization(os.path.join(MEDIA_DIR, "databus_visualization.png"))
    create_memory_blocks(os.path.join(MEDIA_DIR, "memory_blocks.png"))
    create_processor_comparison(os.path.join(MEDIA_DIR, "processor_comparison.png"))
    create_bar_comparison(os.path.join(MEDIA_DIR, "bar_comparison.png"))

    print("\n2. Erstelle PowerPoint-Präsentation...")
    pptx_path = create_presentation()

    print("\n" + "=" * 60)
    print("FERTIG!")
    print(f"Präsentation: {pptx_path}")
    print(f"Medien-Ordner: {MEDIA_DIR}")
    print("=" * 60)

    # Liste der erstellten Dateien
    print("\nErstellte Medien-Dateien:")
    for f in os.listdir(MEDIA_DIR):
        print(f"  - {f}")

if __name__ == "__main__":
    main()
